# ingest_and_chat/nodes.py
"""
Pipeline node functions for the LangGraph RAG ingestion graph.

Each public function is a graph node that receives and returns RAGIngestionState.
All heavy dependencies (LLM, embedding model, genai client) are accessed via
lazy getters from .config to avoid loading at import time.
"""

import os
import sys
import json
import shutil
import datetime
import logging

import psycopg2
from psycopg2.extras import Json
from langchain_text_splitters import RecursiveCharacterTextSplitter, Language

from .states import RAGIngestionState, FileMetadata
from .config import (
    get_llm,
    get_genai_client,
    get_embedding_model,
    OCR_METHOD,
    TRANSCRIPTION_METHOD,
    MAX_BINARY_MB,
    CHUNK_SIZE_TEXT,
    CHUNK_OVERLAP_TEXT,
    CHUNK_SIZE_CODE,
    CHUNK_OVERLAP_CODE,
    LOGS_ROOT,
    TEXT_EXTENSIONS,
    CODE_EXTENSIONS,
    CONFIG_EXTENSIONS,
    PDF_EXTENSIONS,
    OFFICE_EXTENSIONS,
    STRUCTURED_EXTENSIONS,
    IMAGE_EXTENSIONS,
    AUDIO_EXTENSIONS,
    CODE_LANGUAGE_MAP,
    MIME_TYPE_MAP,
)
from .utils import build_project_tree, get_file_stats

logger = logging.getLogger("ingest_and_chat")


# ===========================================================================
# Private Helpers
# ===========================================================================

def _debug_error_with_llm(error_message, context, step_name):
    """Use Gemini to analyze an error and produce a debug summary + log file."""
    llm = get_llm()
    debug_prompt = (
        f"An error occurred during the '{step_name}' step of a RAG ingestion pipeline.\n\n"
        f"--- CONTEXT ---\n"
        f"Target Path: {context.get('target_path', 'N/A')}\n"
        f"PostgreSQL Host: {context.get('pg_host', 'N/A')}\n"
        f"PostgreSQL Database: {context.get('pg_database', 'N/A')}\n"
        f"Last Command: {context.get('last_command', 'N/A')}\n"
        f"Steps Completed: {context.get('steps_completed', [])}\n\n"
        f"--- STDERR ---\n{context.get('last_stderr', 'N/A')}\n\n"
        f"--- ERROR ---\n{error_message}\n\n"
        "Provide:\n1. A clear explanation of the root cause.\n"
        "2. Step-by-step fixes the user can apply.\n"
    )
    try:
        response = llm.invoke(debug_prompt)
        summary = response.content
    except Exception as llm_err:
        summary = f"LLM debug failed: {llm_err}\n\nOriginal error:\n{error_message}"

    log_dir = LOGS_ROOT
    os.makedirs(log_dir, exist_ok=True)
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_path = os.path.join(log_dir, f"error_{step_name}_{timestamp}.txt")

    with open(log_path, "w", encoding="utf-8") as f:
        f.write(f"ERROR LOG -- Step: {step_name}\nTimestamp: {timestamp}\n")
        f.write("=" * 80 + "\n\n--- LLM DEBUG SUMMARY ---\n\n")
        f.write((summary or "N/A") + "\n\n--- RAW ERROR ---\n\n")
        f.write((error_message or "N/A") + "\n")

    return summary, log_path


def _build_context(state):
    return {
        "target_path": state.get("target_path"),
        "pg_host": state.get("pg_host"),
        "pg_database": state.get("pg_database"),
        "last_command": state.get("last_command"),
        "last_stderr": state.get("last_stderr"),
        "steps_completed": state.get("steps_completed"),
    }


def _error_return(state, step_name, error_msg, cmd_str=""):
    """Standard error-state builder shared by every node."""
    logger.error("[%s] Error: %s", step_name, error_msg)
    context = _build_context(state)
    debug_summary, log_path = _debug_error_with_llm(error_msg, context, step_name)
    existing_errors = list(state.get("errors") or [])
    existing_errors.append(error_msg)
    return {
        **state,
        "current_step": step_name,
        "has_error": True,
        "errors": existing_errors,
        "error_log_path": log_path,
        "debug_summary": debug_summary,
        "last_command": cmd_str,
        "last_stderr": error_msg,
    }


def _success_step(state, step_name, updates, cmd_str=""):
    """Standard success-state builder shared by every node."""
    completed = list(state.get("steps_completed") or [])
    completed.append(step_name)
    return {
        **state,
        "current_step": step_name,
        "steps_completed": completed,
        "last_command": cmd_str,
        **updates,
    }


# ---------------------------------------------------------------------------
# Local Output Helpers
# ---------------------------------------------------------------------------

def _relative_path(filepath, target_root):
    abs_file = os.path.abspath(filepath)
    abs_root = os.path.abspath(target_root)
    if os.path.isfile(abs_root):
        return os.path.basename(abs_file)
    try:
        return os.path.relpath(abs_file, abs_root)
    except ValueError:
        return os.path.basename(abs_file)


def _type_to_output_subdir(file_type_or_ext):
    ext = file_type_or_ext.lower().lstrip(".")
    mapping = {
        "py": "code", "js": "code", "ts": "code", "jsx": "code", "tsx": "code",
        "java": "code", "c": "code", "cpp": "code", "h": "code", "hpp": "code",
        "go": "code", "rs": "code", "rb": "code", "php": "code",
        "swift": "code", "kt": "code", "scala": "code", "sh": "code",
        "bash": "code", "r": "code", "sql": "code", "css": "code",
        "scss": "code", "less": "code", "vue": "code", "svelte": "code",
        "md": "text", "txt": "text", "rst": "text", "html": "text",
        "htm": "text", "log": "text", "tex": "text",
        "json": "config", "yaml": "config", "yml": "config", "toml": "config",
        "xml": "config", "env": "config", "ini": "config", "cfg": "config",
        "conf": "config",
        "pdf": "pdf",
        "docx": "office", "pptx": "office",
        "csv": "structured", "tsv": "structured", "xlsx": "structured",
        "xls": "structured",
        "png": "media/images", "jpg": "media/images", "jpeg": "media/images",
        "gif": "media/images", "bmp": "media/images", "tiff": "media/images",
        "webp": "media/images",
        "mp3": "media/audio", "wav": "media/audio", "m4a": "media/audio",
        "ogg": "media/audio", "flac": "media/audio", "aac": "media/audio",
        "wma": "media/audio",
    }
    return mapping.get(ext, "other")


def _save_output_file(output_dir, subdir, rel_path, suffix, content):
    out_path = os.path.join(output_dir, subdir, rel_path + suffix)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(content)
    return out_path


def _save_json(output_dir, subdir, rel_path, suffix, data):
    content = json.dumps(data, indent=2, ensure_ascii=False, default=str)
    return _save_output_file(output_dir, subdir, rel_path, suffix, content)


def _copy_original(output_dir, subdir, rel_path, src_path):
    out_path = os.path.join(output_dir, subdir, rel_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    try:
        shutil.copy2(src_path, out_path)
    except Exception:
        pass
    return out_path


# ---------------------------------------------------------------------------
# File Classification Helper
# ---------------------------------------------------------------------------

def _classify_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    info = {"path": filepath, "extension": ext}

    if ext in CODE_EXTENSIONS:
        info["language"] = CODE_LANGUAGE_MAP.get(ext, "python")
        return "code", info
    if ext in TEXT_EXTENSIONS:
        return "text", info
    if ext in CONFIG_EXTENSIONS:
        return "config", info
    if ext in PDF_EXTENSIONS:
        return "pdf", info
    if ext in OFFICE_EXTENSIONS:
        info["subtype"] = ext.lstrip(".")
        return "office", info
    if ext in STRUCTURED_EXTENSIONS:
        info["subtype"] = ext.lstrip(".")
        return "structured", info
    if ext in IMAGE_EXTENSIONS:
        info["mime"] = MIME_TYPE_MAP.get(ext, "application/octet-stream")
        return "image", info
    if ext in AUDIO_EXTENSIONS:
        info["mime"] = MIME_TYPE_MAP.get(ext, "application/octet-stream")
        return "audio", info
    return "unknown", info


# ---------------------------------------------------------------------------
# Extraction Helpers
# ---------------------------------------------------------------------------

def _extract_pdf(filepath):
    import pymupdf as fitz

    genai_client = get_genai_client()
    doc = fitz.open(filepath)
    full_text = []
    method = "pymupdf"

    for page_num, page in enumerate(doc):
        text = page.get_text("text").strip()
        if len(text) < 50 and page.get_images():
            if genai_client:
                try:
                    pix = page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    from google.genai import types
                    response = genai_client.models.generate_content(
                        model="gemini-2.5-flash",
                        contents=[
                            types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                            "Extract ALL text from this scanned page. Return only the text.",
                        ],
                    )
                    text = response.text.strip()
                    method = "pymupdf+gemini_vision"
                except Exception:
                    pass
        full_text.append(text)

    doc.close()
    return "\n\n".join(full_text), method


def _extract_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs)


def _extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def _ocr_image(filepath, method):
    genai_client = get_genai_client()
    if method == "gemini" and genai_client:
        with open(filepath, "rb") as f:
            img_bytes = f.read()
        ext = os.path.splitext(filepath)[1].lower()
        mime = MIME_TYPE_MAP.get(ext, "image/png")
        from google.genai import types
        response = genai_client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                types.Part.from_bytes(data=img_bytes, mime_type=mime),
                "Extract ALL text visible in this image. If no text is present, "
                "describe the image content in detail. Return only the extracted "
                "or described content.",
            ],
        )
        return response.text.strip(), "gemini_vision"
    else:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return text.strip(), "tesseract"
        except ImportError:
            return "[OCR unavailable -- install pytesseract]", "none"


def _transcribe_audio(filepath, method):
    genai_client = get_genai_client()
    if method == "gemini" and genai_client:
        uploaded = genai_client.files.upload(file=filepath)
        response = genai_client.models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                uploaded,
                "Transcribe this audio completely and accurately. "
                "Return only the transcription text.",
            ],
        )
        return response.text.strip(), "gemini_audio"
    else:
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(filepath)
            return result["text"].strip(), "whisper"
        except ImportError:
            return "[Transcription unavailable -- install openai-whisper]", "none"


def _read_structured(filepath):
    import pandas as pd
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext in (".csv", ".tsv"):
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(filepath, sep=sep, nrows=500)
        elif ext in (".xlsx", ".xls"):
            df = pd.read_excel(filepath, nrows=500)
        else:
            return {"content": "", "column_names": [], "row_count": 0}
        return {
            "content": df.to_csv(index=False),
            "column_names": list(df.columns),
            "row_count": len(df),
            "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        }
    except Exception as e:
        return {"content": "", "column_names": [], "row_count": 0, "error": str(e)}


# ===========================================================================
# NODE 1: Read & Classify Files
# ===========================================================================

def read_and_classify_files(state):
    step_name = "read_and_classify_files"
    target_path = state.get("target_path", "")
    output_dir = state.get("output_dir", "")
    cmd_str = f"Scan & classify: {target_path}"

    logger.info("[%s] Scanning: %s", step_name, target_path)

    try:
        classified = {
            "text": [], "code": [], "config": [], "pdf": [],
            "office": [], "structured": [], "image": [], "audio": [], "unknown": [],
        }

        if os.path.isfile(target_path):
            cat, info = _classify_file(target_path)
            classified[cat].append(info)
        elif os.path.isdir(target_path):
            for root, dirs, files in os.walk(target_path):
                dirs[:] = [
                    d for d in dirs
                    if not d.startswith(".") and d not in
                    {"__pycache__", "node_modules", ".git", "venv", ".venv"}
                ]
                for fname in files:
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(root, fname)
                    cat, info = _classify_file(fpath)
                    classified[cat].append(info)
        else:
            raise FileNotFoundError(f"Path does not exist: {target_path}")

        total = sum(len(v) for v in classified.values())
        if total == 0:
            raise ValueError(f"No files found at {target_path}")

        tree = build_project_tree(target_path)

        if output_dir:
            tree_path = os.path.join(output_dir, "_project_tree.txt")
            with open(tree_path, "w", encoding="utf-8") as f:
                f.write(f"Project Tree for: {target_path}\n")
                f.write("=" * 60 + "\n\n")
                f.write(tree)

            classification_summary = {
                "target_path": target_path,
                "total_files": total,
                "breakdown": {k: len(v) for k, v in classified.items() if v},
                "files_by_category": {
                    k: [{"path": fi["path"], "extension": fi.get("extension", "")}
                        for fi in v]
                    for k, v in classified.items() if v
                },
            }
            cls_path = os.path.join(output_dir, "_classification.json")
            with open(cls_path, "w", encoding="utf-8") as f:
                json.dump(classification_summary, f, indent=2, ensure_ascii=False)

        summary_parts = [f"{k}: {len(v)}" for k, v in classified.items() if v]
        logger.info("[%s] Found %d file(s) -- %s", step_name, total, ", ".join(summary_parts))

        return _success_step(state, step_name, {
            "classified_files": classified,
            "project_tree": tree,
            "last_stdout": f"Classified {total} files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 2: Process Text Documents
# ===========================================================================

def process_text_documents(state):
    step_name = "process_text_documents"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    cmd_str = "Extract text from documents"

    text_categories = ["text", "code", "config", "pdf", "office"]
    file_lists = [classified.get(cat, []) for cat in text_categories]
    all_files = [f for sublist in file_lists for f in sublist]

    if not all_files:
        logger.info("[%s] No text documents to process -- skipping.", step_name)
        return _success_step(state, step_name, {"processed_documents": []}, cmd_str)

    logger.info("[%s] Processing %d text document(s)...", step_name, len(all_files))

    try:
        processed = []
        saved_count = 0

        for finfo in all_files:
            fpath = finfo["path"]
            ext = finfo.get("extension", "")
            extraction_method = "direct_read"
            extracted_text = ""

            try:
                if ext in PDF_EXTENSIONS:
                    extracted_text, extraction_method = _extract_pdf(fpath)
                elif ext == ".docx":
                    extracted_text = _extract_docx(fpath)
                    extraction_method = "python_docx"
                elif ext == ".pptx":
                    extracted_text = _extract_pptx(fpath)
                    extraction_method = "python_pptx"
                else:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        extracted_text = f.read()
            except Exception as read_err:
                extracted_text = f"[extraction error: {read_err}]"
                extraction_method = "error"

            if not extracted_text.strip():
                extracted_text = "[empty file]"

            doc_record = {
                "filepath": fpath,
                "file_type": ext.lstrip("."),
                "language": finfo.get("language"),
                "extracted_text": extracted_text,
                "extraction_method": extraction_method,
            }
            processed.append(doc_record)

            if output_dir:
                rel = _relative_path(fpath, target_path)
                subdir = _type_to_output_subdir(ext)
                _save_output_file(output_dir, subdir, rel, ".extracted.txt", extracted_text)
                saved_count += 1

        logger.info("[%s] Extracted text from %d file(s).", step_name, len(processed))

        return _success_step(state, step_name, {
            "processed_documents": processed,
            "last_stdout": f"Processed {len(processed)} text documents",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 3: Process Media Files
# ===========================================================================

def process_media_files(state):
    step_name = "process_media_files"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    cmd_str = "OCR images / transcribe audio"

    images = classified.get("image", [])
    audios = classified.get("audio", [])

    if not images and not audios:
        logger.info("[%s] No media files to process -- skipping.", step_name)
        return _success_step(state, step_name, {"processed_media": []}, cmd_str)

    logger.info("[%s] Processing %d image(s) + %d audio file(s)...",
                step_name, len(images), len(audios))

    try:
        processed = []

        for img_info in images:
            fpath = img_info["path"]
            try:
                transcript, method = _ocr_image(fpath, OCR_METHOD)
            except Exception as e:
                transcript, method = f"[OCR failed: {e}]", "error"

            processed.append({
                "filepath": fpath, "file_type": "image",
                "transcript": transcript, "extraction_method": method,
            })
            logger.info("  OCR (%s): %s", method, os.path.basename(fpath))

            if output_dir:
                rel = _relative_path(fpath, target_path)
                _save_output_file(output_dir, "media/images", rel, ".ocr.txt", transcript)
                _copy_original(output_dir, "media/images", rel, fpath)

        for aud_info in audios:
            fpath = aud_info["path"]
            try:
                transcript, method = _transcribe_audio(fpath, TRANSCRIPTION_METHOD)
            except Exception as e:
                transcript, method = f"[Transcription failed: {e}]", "error"

            processed.append({
                "filepath": fpath, "file_type": "audio",
                "transcript": transcript, "extraction_method": method,
            })
            logger.info("  Transcribed (%s): %s", method, os.path.basename(fpath))

            if output_dir:
                rel = _relative_path(fpath, target_path)
                _save_output_file(output_dir, "media/audio", rel, ".transcript.txt", transcript)
                _copy_original(output_dir, "media/audio", rel, fpath)

        logger.info("[%s] Processed %d media file(s).", step_name, len(processed))

        return _success_step(state, step_name, {
            "processed_media": processed,
            "last_stdout": f"Processed {len(processed)} media files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 4: Process Structured Files
# ===========================================================================

def process_structured_files(state):
    step_name = "process_structured_files"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    structured = classified.get("structured", [])
    cmd_str = "Read structured files"

    if not structured:
        logger.info("[%s] No structured files -- skipping.", step_name)
        return _success_step(state, step_name, {"processed_structured": []}, cmd_str)

    logger.info("[%s] Reading %d structured file(s)...", step_name, len(structured))

    try:
        processed = []
        for finfo in structured:
            fpath = finfo["path"]
            result = _read_structured(fpath)
            record = {"filepath": fpath, "file_type": finfo.get("subtype", "csv"), **result}
            processed.append(record)
            logger.info("  Read: %s -- %d rows, %d columns",
                        os.path.basename(fpath),
                        result.get("row_count", 0),
                        len(result.get("column_names", [])))

            if output_dir:
                rel = _relative_path(fpath, target_path)
                preview = {
                    "filepath": fpath,
                    "file_type": finfo.get("subtype", "csv"),
                    "column_names": result.get("column_names", []),
                    "dtypes": result.get("dtypes", {}),
                    "row_count": result.get("row_count", 0),
                    "first_10_rows": result.get("content", "").split("\n")[:11],
                }
                _save_json(output_dir, "structured", rel, ".preview.json", preview)
                _copy_original(output_dir, "structured", rel, fpath)

        logger.info("[%s] Processed %d structured file(s).", step_name, len(processed))
        return _success_step(state, step_name, {
            "processed_structured": processed,
            "last_stdout": f"Processed {len(processed)} structured files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 5: Generate Metadata
# ===========================================================================

def generate_metadata(state):
    step_name = "generate_metadata"
    cmd_str = "Gemini structured metadata generation"
    project_tree = state.get("project_tree", "N/A")
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")

    docs = state.get("processed_documents") or []
    media = state.get("processed_media") or []
    structured = state.get("processed_structured") or []

    all_items = (
        [(d["filepath"], d.get("extracted_text", "")[:6000], d.get("file_type", ""))
         for d in docs]
        + [(m["filepath"], m.get("transcript", "")[:6000], m.get("file_type", ""))
           for m in media]
        + [(s["filepath"], s.get("content", "")[:3000], s.get("file_type", ""))
           for s in structured]
    )

    if not all_items:
        logger.info("[%s] Nothing to generate metadata for.", step_name)
        return _success_step(state, step_name, {"file_metadata": {}}, cmd_str)

    logger.info("[%s] Generating metadata for %d file(s)...", step_name, len(all_items))

    try:
        llm = get_llm()
        metadata_map = {}
        structured_llm = llm.with_structured_output(FileMetadata)

        for filepath, content_sample, file_type in all_items:
            stats = get_file_stats(filepath)
            prompt = (
                "You are an expert data cataloger. Generate comprehensive metadata "
                "for this file.\n\n"
                f"FILE PATH: {filepath}\n"
                f"FILE TYPE: {file_type}\n"
                f"FILE SIZE: {stats.get('file_size_bytes', 0)} bytes\n"
                f"CREATED: {stats.get('file_created_date', 'unknown')}\n"
                f"MODIFIED: {stats.get('file_modified_date', 'unknown')}\n\n"
                f"PROJECT TREE:\n{project_tree[:2000]}\n\n"
                f"CONTENT SAMPLE:\n{content_sample}\n"
            )

            try:
                llm_meta = structured_llm.invoke(prompt)
            except Exception as e:
                llm_meta = {
                    "summary": "Metadata extraction failed.",
                    "topics": [],
                    "key_entities": [],
                    "content_category": "unknown",
                    "quality_notes": f"API Error: {str(e)}",
                }

            full_meta = {
                **llm_meta, **stats,
                "source_file_type": file_type,
                "source_filepath": filepath,
                "project_tree_snippet": project_tree[:500],
            }
            metadata_map[filepath] = full_meta
            logger.info("  %s: %s", os.path.basename(filepath),
                        llm_meta.get("content_category", "?"))

            if output_dir:
                rel = _relative_path(filepath, target_path)
                ext = os.path.splitext(filepath)[1].lower()
                subdir = _type_to_output_subdir(ext)
                _save_json(output_dir, subdir, rel, ".metadata.json", full_meta)

        if output_dir and metadata_map:
            combined_path = os.path.join(output_dir, "metadata", "all_metadata.json")
            os.makedirs(os.path.dirname(combined_path), exist_ok=True)
            with open(combined_path, "w", encoding="utf-8") as f:
                json.dump(metadata_map, f, indent=2, ensure_ascii=False, default=str)

        logger.info("[%s] Generated metadata for %d file(s).", step_name, len(metadata_map))

        return _success_step(state, step_name, {
            "file_metadata": metadata_map,
            "last_stdout": f"Generated metadata for {len(metadata_map)} files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 6: Setup PostgreSQL
# ===========================================================================

_REQUIRED_TABLES = {
    "document_chunks": {
        "columns": {
            "id", "filepath", "file_type", "chunk_index", "total_chunks",
            "content", "metadata", "embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE document_chunks (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                chunk_index INTEGER NOT NULL,
                total_chunks INTEGER NOT NULL,
                content TEXT NOT NULL,
                metadata JSONB,
                embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
    "structured_files": {
        "columns": {
            "id", "filepath", "file_type", "content", "column_names",
            "row_count", "metadata", "summary_embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE structured_files (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                content TEXT NOT NULL,
                column_names JSONB,
                row_count INTEGER,
                metadata JSONB,
                summary_embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
    "media_files": {
        "columns": {
            "id", "filepath", "file_type", "binary_data", "transcript",
            "metadata", "transcript_embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE media_files (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                binary_data BYTEA,
                transcript TEXT,
                metadata JSONB,
                transcript_embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
}


def _get_existing_columns(cur, table_name):
    cur.execute("""
        SELECT column_name
        FROM information_schema.columns
        WHERE table_schema = 'public' AND table_name = %s
    """, (table_name,))
    return {row[0] for row in cur.fetchall()}


def setup_postgres(state):
    step_name = "setup_postgres"
    cmd_str = "CREATE TABLES (document_chunks, structured_files, media_files)"
    logger.info("[%s] Setting up PostgreSQL tables...", step_name)

    try:
        conn = psycopg2.connect(
            host=state.get("pg_host"),
            port=state.get("pg_port"),
            user=state.get("pg_username"),
            password=state.get("pg_password"),
            dbname=state.get("pg_database"),
            connect_timeout=10,
        )
        conn.autocommit = True
        cur = conn.cursor()

        cur.execute("CREATE EXTENSION IF NOT EXISTS vector;")
        logger.info("  pgvector extension enabled")

        for table_name, spec in _REQUIRED_TABLES.items():
            existing_cols = _get_existing_columns(cur, table_name)
            if not existing_cols:
                cur.execute(spec["create_sql"])
                logger.info("  Created table: %s", table_name)
            elif not spec["columns"].issubset(existing_cols):
                missing = spec["columns"] - existing_cols
                logger.warning("  Table '%s' stale schema (missing: %s) -- recreating",
                               table_name, missing)
                cur.execute(f"DROP TABLE IF EXISTS {table_name} CASCADE;")
                cur.execute(spec["create_sql"])
                logger.info("  Recreated table: %s", table_name)
            else:
                logger.info("  Table exists with correct schema: %s", table_name)

        cur.close()
        conn.close()
        logger.info("[%s] Tables ready.", step_name)

        return _success_step(state, step_name, {
            "last_stdout": "PostgreSQL tables created/verified",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 7: Vectorize & Store
# ===========================================================================

def _get_text_splitter(language):
    if language:
        try:
            lang_enum = Language(language)
            return RecursiveCharacterTextSplitter.from_language(
                language=lang_enum,
                chunk_size=CHUNK_SIZE_CODE,
                chunk_overlap=CHUNK_OVERLAP_CODE,
            )
        except ValueError:
            pass
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE_TEXT,
        chunk_overlap=CHUNK_OVERLAP_TEXT,
        separators=["\n\n", "\n", ". ", " ", ""],
    )


def vectorize_and_store(state):
    step_name = "vectorize_and_store"
    cmd_str = "Vectorize & INSERT into PostgreSQL"
    metadata_map = state.get("file_metadata") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    embedding_model = get_embedding_model()

    docs = state.get("processed_documents") or []
    media = state.get("processed_media") or []
    structured = state.get("processed_structured") or []

    logger.info("[%s] Vectorizing: docs=%d  media=%d  structured=%d",
                step_name, len(docs), len(media), len(structured))

    try:
        conn = psycopg2.connect(
            host=state.get("pg_host"),
            port=state.get("pg_port"),
            user=state.get("pg_username"),
            password=state.get("pg_password"),
            dbname=state.get("pg_database"),
        )
        cur = conn.cursor()
        total_inserted = 0
        saved_chunks = 0

        # ----- 1. Document Chunks -----
        for doc in docs:
            filepath = doc["filepath"]
            text = doc.get("extracted_text", "")
            language = doc.get("language")
            meta = metadata_map.get(filepath, {})

            if not text or text.startswith("["):
                continue

            splitter = _get_text_splitter(language)
            chunks = splitter.split_text(text)
            if not chunks:
                continue

            embeddings = embedding_model.encode(chunks).tolist()
            local_chunks_data = []

            for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                chunk_meta = {**meta, "chunk_index": idx, "total_chunks": len(chunks)}
                formatted_emb = f"[{','.join(map(str, emb))}]"
                cur.execute("""
                    INSERT INTO document_chunks
                        (filepath, file_type, chunk_index, total_chunks,
                         content, metadata, embedding)
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                """, (
                    filepath, doc.get("file_type", "unknown"),
                    idx, len(chunks), chunk, Json(chunk_meta), formatted_emb,
                ))
                total_inserted += 1
                local_chunks_data.append({
                    "chunk_index": idx, "total_chunks": len(chunks),
                    "char_count": len(chunk), "content": chunk,
                })

            logger.info("  %s: %d chunk(s)", os.path.basename(filepath), len(chunks))

            if output_dir:
                rel = _relative_path(filepath, target_path)
                ext = os.path.splitext(filepath)[1].lower()
                subdir = _type_to_output_subdir(ext)
                _save_json(output_dir, subdir, rel, ".chunks.json", {
                    "filepath": filepath,
                    "file_type": doc.get("file_type", "unknown"),
                    "language": language,
                    "chunk_config": {
                        "chunk_size": CHUNK_SIZE_CODE if language else CHUNK_SIZE_TEXT,
                        "chunk_overlap": CHUNK_OVERLAP_CODE if language else CHUNK_OVERLAP_TEXT,
                        "splitter": f"code:{language}" if language else "recursive_text",
                    },
                    "total_chunks": len(chunks),
                    "chunks": local_chunks_data,
                })
                saved_chunks += 1

        # ----- 2. Media Files -----
        for item in media:
            filepath = item["filepath"]
            transcript = item.get("transcript", "")
            meta = metadata_map.get(filepath, {})

            binary_data = None
            try:
                file_size_mb = os.path.getsize(filepath) / (1024 * 1024)
                if file_size_mb <= MAX_BINARY_MB:
                    with open(filepath, "rb") as f:
                        binary_data = f.read()
                else:
                    meta["binary_stored"] = False
                    meta["binary_note"] = (
                        f"File too large ({file_size_mb:.1f}MB > {MAX_BINARY_MB}MB limit)."
                    )
            except Exception:
                pass

            transcript_emb = None
            if transcript and not transcript.startswith("["):
                emb = embedding_model.encode(transcript).tolist()
                transcript_emb = f"[{','.join(map(str, emb))}]"

            cur.execute("""
                INSERT INTO media_files
                    (filepath, file_type, binary_data, transcript,
                     metadata, transcript_embedding)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (
                filepath, item.get("file_type", "unknown"),
                psycopg2.Binary(binary_data) if binary_data else None,
                transcript, Json(meta), transcript_emb,
            ))
            total_inserted += 1

        # ----- 3. Structured Files -----
        for item in structured:
            filepath = item["filepath"]
            content = item.get("content", "")
            meta = metadata_map.get(filepath, {})
            meta["column_names"] = item.get("column_names", [])
            meta["row_count"] = item.get("row_count", 0)
            meta["dtypes"] = item.get("dtypes", {})

            summary_text = (
                f"File: {os.path.basename(filepath)}. "
                f"Columns: {', '.join(item.get('column_names', []))}. "
                f"Rows: {item.get('row_count', 0)}. "
                f"{meta.get('summary', '')}"
            )
            emb = embedding_model.encode(summary_text).tolist()
            summary_emb = f"[{','.join(map(str, emb))}]"

            cur.execute("""
                INSERT INTO structured_files
                    (filepath, file_type, content, column_names,
                     row_count, metadata, summary_embedding)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            """, (
                filepath, item.get("file_type", "csv"),
                content, Json(item.get("column_names", [])),
                item.get("row_count", 0), Json(meta), summary_emb,
            ))
            total_inserted += 1

        conn.commit()
        cur.close()
        conn.close()

        # === SAVE LOCALLY: final overview ===
        if output_dir:
            overview = {
                "session_id": state.get("session_id"),
                "target_path": target_path,
                "ingestion_timestamp": datetime.datetime.now().isoformat(),
                "total_records_inserted": total_inserted,
                "steps_completed": list(state.get("steps_completed") or []) + [step_name],
                "summary": {
                    "documents_processed": len(docs),
                    "media_processed": len(media),
                    "structured_processed": len(structured),
                    "total_chunks_created": saved_chunks,
                },
                "postgresql": {
                    "host": state.get("pg_host"),
                    "port": state.get("pg_port"),
                    "database": state.get("pg_database"),
                    "tables": ["document_chunks", "media_files", "structured_files"],
                },
                "project_tree": state.get("project_tree", ""),
            }
            overview_path = os.path.join(output_dir, "_overview.json")
            with open(overview_path, "w", encoding="utf-8") as f:
                json.dump(overview, f, indent=2, ensure_ascii=False, default=str)

            readme_path = os.path.join(output_dir, "_README.txt")
            with open(readme_path, "w", encoding="utf-8") as f:
                f.write("=" * 60 + "\n")
                f.write("  RAG Pipeline -- Processed Output\n")
                f.write("=" * 60 + "\n\n")
                f.write(f"Source:     {target_path}\n")
                f.write(f"Timestamp:  {overview['ingestion_timestamp']}\n")
                f.write(f"Records:    {total_inserted} inserted into PostgreSQL\n\n")
                f.write("FOLDER STRUCTURE:\n" + "-" * 40 + "\n")
                f.write("  code/          Source code text, chunks, metadata\n")
                f.write("  text/          Text/doc text, chunks, metadata\n")
                f.write("  config/        Config file text, chunks, metadata\n")
                f.write("  pdf/           PDF text, chunks, metadata\n")
                f.write("  office/        DOCX/PPTX text, chunks, metadata\n")
                f.write("  media/images/  Original images + OCR (.ocr.txt)\n")
                f.write("  media/audio/   Original audio + transcripts (.transcript.txt)\n")
                f.write("  structured/    Original CSV/XLSX + previews (.preview.json)\n")
                f.write("  metadata/      Combined metadata for ALL files\n")

        logger.info("[%s] Stored %d record(s) in PostgreSQL.", step_name, total_inserted)

        return _success_step(state, step_name, {
            "records_inserted": total_inserted,
            "last_stdout": f"Inserted {total_inserted} records",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)
